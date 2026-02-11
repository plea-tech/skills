#!/usr/bin/env python3
"""当 run_markitdown 不可用时，用 pypdf + python-docx 将 PDF/DOCX 转为纯文本到 text/。"""
import json
import re
import sys
from pathlib import Path

def safe_md_name(filename: str) -> str:
    stem = Path(filename).stem
    name = (re.sub(r"[\s\-]+", "_", stem)).strip() or "doc"
    return name + ".md"

def extract_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages):
        t = page.extract_text()
        if t:
            parts.append(f"## Page {i+1}\n\n{t}")
    return "\n\n".join(parts) if parts else ""

def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for p in doc.paragraphs:
        parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n\n".join(p for p in parts if p.strip())

def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides):
        part = [f"## Slide {i+1}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                part.append(shape.text)
        parts.append("\n\n".join(part))
    return "\n\n".join(parts)

def main():
    worklist_path = Path(sys.argv[1]) if len(sys.argv) > 1 else None
    if not worklist_path or not worklist_path.is_file():
        print("Usage: extract_text_fallback.py <worklist.json>", file=sys.stderr)
        sys.exit(1)
    with open(worklist_path, encoding="utf-8") as f:
        wl = json.load(f)
    extracted_root = Path(wl["extracted_root"])
    text_dir = extracted_root / "text"
    text_dir.mkdir(parents=True, exist_ok=True)
    safe_name = lambda f: Path(f).stem.replace(" ", "_").replace("-", "_") + ".md"
    for doc in wl.get("documents", []):
        path = Path(doc["path"])
        if not path.is_file():
            continue
        out = text_dir / safe_name(doc.get("filename", path.name))
        try:
            if path.suffix.lower() == ".pdf":
                text = extract_pdf(path)
            elif path.suffix.lower() in (".docx", ".doc"):
                text = extract_docx(path)
            elif path.suffix.lower() == ".pptx":
                text = extract_pptx(path)
            else:
                continue
            out.write_text(text or "(无文本)", encoding="utf-8")
            print("OK", out.name)
        except Exception as e:
            print("FAIL", path.name, str(e), file=sys.stderr)

if __name__ == "__main__":
    main()
